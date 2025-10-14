import streamlit as st
import io
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from .mkcsv import mkcsv_gui
from .mkpptx import mkpptx0_gui

def mkslide_gui():

    # ファイルアップロード
    uploaded_file = st.file_uploader("Excel/CSVファイルをアップロード", type=["xlsx", "xls", "csv"])
    option_form = st.radio("MKSLIDEが作成したCSVファイルですか？", ["Yes", "No"], index = 1, horizontal = True)
    
    if uploaded_file and option_form == "No":
       mkcsv_gui(uploaded_file)
    elif uploaded_file and option_form == "Yes":
       df = pd.read_csv(uploaded_file)
       st.dataframe(df)
    
    uploaded_pict = st.file_uploader("画像ファイルを選択してください（複数可）",
                    type=["png", "jpg", "jpeg"], accept_multiple_files=True)
    disply_form = st.radio("アップロードした画像を表示しますか？", ["Yes", "No"], index = 1, horizontal = True)
    
    if uploaded_pict:
       st.success(f"{len(uploaded_pict)} 件の画像をアップロードしました")
       images = {pic.name: Image.open(pic) for pic in uploaded_pict}
       if disply_form == "Yes":
          cols = st.columns(3)
          for i, (name, image) in enumerate(images.items()):
              col = cols[i % 3]
              col.image(image, caption=name, use_container_width=True)
              
       selected_exam = st.selectbox("試験を選んでください", ["指定しない"]+df["試験"].unique().tolist())
       selected_face = st.selectbox("測定面を選んでください", ["指定しない"]+df["測定面"].unique().tolist())
       selected_cath = st.selectbox("正極を選んでください", ["指定しない"]+df["正極"].unique().tolist()) 
       selected_mesu = st.selectbox("測定を選んでください", ["指定しない"]+df["測定"].unique().tolist())
       selected_elec = st.selectbox("電解液を選んでください", ["指定しない"]+df["電解液"].unique().tolist())
       selected_magn = st.selectbox("倍率を選んでください", ["指定しない"]+df["倍率"].unique().tolist())

       condition = pd.Series(True, index=df.index)
       if selected_exam != "指定しない":
          condition &= (df["試験"] == selected_exam)
       if selected_face != "指定しない":
          condition &= (df["測定面"] == selected_face)
       if selected_cath != "指定しない":
          condition &= (df["正極"] == selected_cath)
       if selected_mesu != "指定しない":
          condition &= (df["測定"] == selected_mesu)
       if selected_elec != "指定しない":
          condition &= (df["電解液"] == selected_elec)
       if selected_magn != "指定しない":
          condition &= (df["倍率"] == selected_magn)

       result = df.loc[condition, "ファイル名"]
       
       # 実際にアップロードされた画像だけを残す
       st.write(result)
       result = result.tolist()
       result = [name for name in result if name in images]
        
       if len(result) == 0:
          st.warning("該当する画像がありません。")
       else:
          result_cols = st.columns(3)
          for i, name in enumerate(result):
              if name in images:
                 image = images[name]
                 col_result = result_cols[i % 3]
                 col_result.image(image, caption=name, use_container_width=True)

       if len(result) > 0:
          prs = Presentation()
          for name in result:
              if name in images:
                 image = images[name]

                 tmp_path = f"tmp_{name}"
                 image.save(tmp_path)

                 slide = prs.slides.add_slide(prs.slide_layouts[5]) 
                 slide.shapes.add_picture(tmp_path, Inches(1), Inches(1), width=Inches(6))

          pptx_path = "output.pptx"
          prs.save(pptx_path)
          st.success(f"PPTXに画像を追加しました: {pptx_path}")

          with open(pptx_path, "rb") as f:
               st.download_button("PPTXをダウンロード", f, file_name="output.pptx")

# MODULE ERROR MESSAGE
if __name__ == "__main__":
   raise RuntimeError("Do not run this file directly; use it as a module.")
