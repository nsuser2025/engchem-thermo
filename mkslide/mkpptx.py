import streamlit as st
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt

def mkpptx_gui(df, images, result):
    uploaded_template = st.file_uploader("PPTXテンプレートをアップロード", type=["pptx"])
    if uploaded_template is None:
        st.warning("PPTXテンプレートをアップロードしてください。")
        return
    template_bytes = BytesIO(uploaded_template.read())
    prs = Presentation(template_bytes)
    
    cover_slide = prs.slides[0]

    rows, cols_num = 2, 3
    width = Inches(2.4) 
    height = Inches(1.8)
    spacing_x, spacing_y = Inches(0.2), Inches(0.2)
    left_margin, bottom_margin = Inches(0.5), Inches(0.5)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    #total_height = rows * height + (rows - 1) * spacing_y
    max_per_slide = rows * cols_num

    for batch_start in range(0, len(result), max_per_slide):
        batch = result[batch_start: batch_start + max_per_slide]
        slide_layout_index = 2
        slide_layout = prs.slide_layouts[slide_layout_index]
        image_slide = prs.slides.add_slide(slide_layout)

        total_height = rows * height + (rows - 1) * spacing_y

        for idx, name in enumerate(batch):
            if name in images:
               image = images[name]
               tmp_path = f"tmp_{name}"
               image.save(tmp_path)
               row = idx // cols_num
               col = idx % cols_num
               left = left_margin + col * (width + spacing_x)
               top = slide_height - total_height - bottom_margin + row * (height + spacing_y)
               image_slide.shapes.add_picture(tmp_path, left, top, width=width, height=height) 
            
            condition_row = df[df["ファイル名"] == name]
            if not condition_row.empty:
               condition_row = condition_row.iloc[0]
               cond_text = (f"{condition_row['試験']}-"
                            f"{condition_row['測定面']}-"
                            f"{condition_row['正極']}-\n"
                            f"{condition_row['測定']}-"
                            f"{condition_row['電解液']}-"
                            f"{condition_row['倍率']}")

               textbox_left = left
               textbox_top = top + height + Inches(0.05)  # 画像下に少し余白
               textbox_width = width
               textbox_height = Inches(0.35)
               textbox = image_slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
               text_frame = textbox.text_frame
               text_frame.clear()
                
               text_frame.margin_left = Pt(2)
               text_frame.margin_right = Pt(2)
               text_frame.margin_top = Pt(2)
               text_frame.margin_bottom = Pt(2) 

               for line in cond_text.split("\n"):
                   p = text_frame.add_paragraph()
                   p.text = line
                   p.font.size = Pt(9)

    pptx_path = "output.pptx"
    prs.save(pptx_path)
    st.success(f"PPTXに画像を追加しました: {pptx_path}")

    with open(pptx_path, "rb") as f:
         st.download_button("PPTXをダウンロード", f, file_name="output.pptx")

# MODULE ERROR MESSAGE
if __name__ == "__main__":
   raise RuntimeError("Do not run this file directly; use it as a module.")
